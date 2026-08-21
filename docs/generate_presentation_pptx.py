#!/usr/bin/env python3
"""PPTX-слайды Альфа-будущее — те же 4 слайда, что в PDF."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "slides" / "alfa-budushchee-slides-market-6-7.pptx"

# Widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

RED = RGBColor(0xEF, 0x31, 0x24)
INK = RGBColor(0x11, 0x11, 0x11)
MUTED = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF6, 0xF6, 0xF6)
SOFT = RGBColor(0xFF, 0xF5, 0xF4)
LINE = RGBColor(0xE8, 0xE8, 0xE8)


def set_run(p, text: str, size: int, bold: bool = False, color=INK):
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return run


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p, text, size, bold, color)
    return box


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    # softer corners
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_header(slide, title: str, page: int, total: int = 4):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    add_textbox(slide, Inches(0.4), Inches(0.12), Inches(4), Inches(0.35), "АЛЬФА · БУДУЩЕЕ", 12, True, WHITE)
    add_textbox(
        slide,
        SLIDE_W - Inches(1.2),
        Inches(0.12),
        Inches(0.9),
        Inches(0.35),
        f"{page}/{total}",
        11,
        False,
        WHITE,
        PP_ALIGN.RIGHT,
    )
    add_textbox(slide, Inches(0.4), Inches(0.7), Inches(12), Inches(0.5), title, 28, True, INK)


def add_footer(slide):
    add_textbox(
        slide,
        Inches(0.4),
        SLIDE_H - Inches(0.35),
        Inches(8),
        Inches(0.25),
        "Фича банка · не стартап · авг 2026",
        9,
        False,
        MUTED,
    )


def card_with_text(slide, left, top, w, h, fill, lines: list[tuple[str, int, bool, object]], align=PP_ALIGN.CENTER):
    add_rect(slide, left, top, w, h, fill, LINE)
    # stack text roughly centered
    n = len(lines)
    block_h = sum(Pt(sz + 6).pt for _, sz, _, _ in lines)  # approx
    # simpler: fixed offsets
    y = top + Inches(0.2)
    for text, size, bold, color in lines:
        add_textbox(slide, left + Inches(0.1), y, w - Inches(0.2), Inches(0.4), text, size, bold, color, align)
        y += Inches(0.28) if size < 18 else Inches(0.38)


# ─── slides ──────────────────────────────────────────────────────────────────


def slide_market(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_header(slide, "Рынок · 2026", 1)

    trends = [
        ("01", "Дешевле", "ETF > активные фонды"),
        ("02", "Шире доступ", "Private → retail"),
        ("03", "AI", "Персонализация"),
        ("04", "CX", "Дистрибуция > продукт"),
    ]
    tw = Inches(2.9)
    gap = Inches(0.2)
    left0 = Inches(0.4)
    top = Inches(1.4)
    for i, (num, title, sub) in enumerate(trends):
        x = left0 + i * (tw + gap)
        add_rect(slide, x, top, tw, Inches(2.0), WHITE, LINE)
        add_textbox(slide, x + Inches(0.2), top + Inches(0.25), tw - Inches(0.3), Inches(0.5), num, 28, True, RED)
        add_textbox(slide, x + Inches(0.2), top + Inches(0.85), tw - Inches(0.3), Inches(0.4), title, 18, True, INK)
        add_textbox(slide, x + Inches(0.2), top + Inches(1.3), tw - Inches(0.3), Inches(0.5), sub, 12, False, MUTED)

    stats = [
        ("$943B", "model portfolios", BG),
        ("$308B", "BlackRock · ~⅓", BG),
        ("→ CX+AI", "где конкурировать", SOFT),
    ]
    sw = Inches(3.95)
    stop = Inches(3.7)
    for i, (v, l, fill) in enumerate(stats):
        x = left0 + i * (sw + gap)
        add_rect(slide, x, stop, sw, Inches(1.5), fill, LINE)
        add_textbox(slide, x, stop + Inches(0.35), sw, Inches(0.55), v, 28, True, RED if i == 2 else INK, PP_ALIGN.CENTER)
        add_textbox(slide, x, stop + Inches(0.95), sw, Inches(0.35), l, 12, False, MUTED, PP_ALIGN.CENTER)

    add_textbox(
        slide,
        Inches(0.4),
        Inches(5.6),
        Inches(12.5),
        Inches(0.4),
        "Выигрывает не продукт, а путь клиента + объяснение",
        16,
        True,
        INK,
        PP_ALIGN.CENTER,
    )
    add_footer(slide)


def slide_competitors(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Конкуренты → наша ниша", 2)

    # Left panel
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(5.8), Inches(4.6), WHITE, LINE)
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5), Inches(0.35), "Где сильные брокеры", 14, True, INK)

    axes = [
        ("Каталог / терминал", 5),
        ("Обучение рынку", 4),
        ("Робот-портфель", 4),
        ("Первый шаг 18–26", 2),
        ("Банк → взнос", 2),
        ("ИИ объясняет выбор", 1),
    ]
    bar_left = Inches(3.4)
    bar_max = Inches(2.4)
    for i, (label, n) in enumerate(axes):
        y = Inches(2.05) + i * Inches(0.55)
        add_textbox(slide, Inches(0.6), y, Inches(2.7), Inches(0.3), label, 11, False, MUTED)
        # bg bar
        add_rect(slide, bar_left, y + Inches(0.05), bar_max, Inches(0.2), BG, None)
        # fill
        fill_w = int(bar_max * n / 5)
        add_rect(slide, bar_left, y + Inches(0.05), fill_w, Inches(0.2), RED, None)

    # Right panel
    add_rect(slide, Inches(6.5), Inches(1.35), Inches(6.4), Inches(4.6), SOFT, LINE)
    add_textbox(slide, Inches(6.7), Inches(1.5), Inches(5), Inches(0.35), "Белое пятно", 14, True, RED)

    steps = ["кэшбэк / остаток", "цель + квиз", "LQDT + «почему»", "взнос от 100 ₽"]
    for i, s in enumerate(steps):
        y = Inches(2.1) + i * Inches(0.85)
        add_rect(slide, Inches(7.3), y, Inches(4.8), Inches(0.55), WHITE, LINE)
        add_textbox(
            slide,
            Inches(7.3),
            y + Inches(0.1),
            Inches(4.8),
            Inches(0.4),
            f"{i + 1}.  {s}",
            14,
            True,
            INK,
            PP_ALIGN.CENTER,
        )
        if i < len(steps) - 1:
            add_textbox(
                slide,
                Inches(7.3),
                y + Inches(0.52),
                Inches(4.8),
                Inches(0.3),
                "↓",
                14,
                True,
                RED,
                PP_ALIGN.CENTER,
            )

    add_textbox(
        slide,
        Inches(0.4),
        Inches(6.15),
        Inches(12.5),
        Inches(0.35),
        "Они продают доступ к рынку · мы — безопасный первый шаг",
        14,
        True,
        INK,
        PP_ALIGN.CENTER,
    )
    add_footer(slide)


def slide_personalization(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "6 · Персонализация", 3)

    boxes = [
        ("ДАННЫЕ", "остаток\nкэшбэк\nцель", BG),
        ("ПРАВИЛА", "18+\nwhitelist\nconservative", SOFT),
        ("LLM", "только\n«почему»\n+ fallback", BG),
    ]
    bw = Inches(2.8)
    gap = Inches(0.9)
    total = 3 * bw + 2 * gap
    x0 = (SLIDE_W - total) / 2
    top = Inches(1.5)

    for i, (title, body, fill) in enumerate(boxes):
        x = x0 + i * (bw + gap)
        add_rect(slide, x, top, bw, Inches(2.6), fill, LINE)
        add_textbox(slide, x, top + Inches(0.25), bw, Inches(0.4), title, 14, True, RED, PP_ALIGN.CENTER)
        for j, line in enumerate(body.split("\n")):
            add_textbox(
                slide,
                x,
                top + Inches(0.85) + j * Inches(0.4),
                bw,
                Inches(0.35),
                line,
                16,
                True,
                INK,
                PP_ALIGN.CENTER,
            )
        if i < 2:
            # arrow between
            ax = x + bw + Inches(0.15)
            add_textbox(slide, ax, top + Inches(1.05), Inches(0.6), Inches(0.4), "→", 28, True, RED, PP_ALIGN.CENTER)

    # Product card
    add_rect(slide, Inches(3.5), Inches(4.5), Inches(6.3), Inches(0.7), INK, None)
    add_textbox(
        slide,
        Inches(3.5),
        Inches(4.6),
        Inches(6.3),
        Inches(0.5),
        "→  Product Card: LQDT + 2 альтернативы",
        16,
        True,
        WHITE,
        PP_ALIGN.CENTER,
    )

    bans = ["✕  не выбирает тикер", "✕  не обещает %", "✕  не ставит ордер"]
    bw2 = Inches(3.5)
    bx0 = Inches(1.2)
    for i, b in enumerate(bans):
        x = bx0 + i * (bw2 + Inches(0.3))
        add_rect(slide, x, Inches(5.5), bw2, Inches(0.55), WHITE, LINE)
        add_textbox(slide, x, Inches(5.58), bw2, Inches(0.4), b, 12, True, RED, PP_ALIGN.CENTER)

    add_footer(slide)


def slide_finance(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "7 · Финмодель фичи", 4)

    add_textbox(
        slide,
        Inches(0.4),
        Inches(1.25),
        Inches(12.5),
        Inches(0.35),
        "700 ₽ чек ≠ бизнес   ·   кейс = 3 слоя",
        16,
        True,
        RED,
        PP_ALIGN.CENTER,
    )

    layers = [
        ("ЛИФТ", "11,1%", "+10% счетов", SOFT),
        ("MIX УК", "30%", "AKMM после LQDT", WHITE),
        ("PRIMARY", "94 млн ₽", "удержание 3г", WHITE),
    ]
    lw = Inches(3.95)
    gap = Inches(0.25)
    for i, (t, v, s, fill) in enumerate(layers):
        x = Inches(0.4) + i * (lw + gap)
        add_rect(slide, x, Inches(1.75), lw, Inches(2.2), fill, LINE)
        add_textbox(slide, x, Inches(1.95), lw, Inches(0.35), t, 12, True, MUTED, PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.45), lw, Inches(0.6), v, 32, True, RED, PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(3.3), lw, Inches(0.35), s, 13, False, INK, PP_ALIGN.CENTER)

    # Flow
    flow = [("цель", False), ("квиз", False), ("LQDT", True), ("AKMM", True)]
    fw = Inches(2.0)
    fx0 = Inches(0.4)
    for i, (label, hot) in enumerate(flow):
        x = fx0 + i * (fw + Inches(0.55))
        add_rect(slide, x, Inches(4.3), fw, Inches(0.7), RED if hot else BG, None if hot else LINE)
        add_textbox(
            slide,
            x,
            Inches(4.42),
            fw,
            Inches(0.45),
            label,
            14,
            True,
            WHITE if hot else INK,
            PP_ALIGN.CENTER,
        )
        if i < len(flow) - 1:
            add_textbox(
                slide,
                x + fw,
                Inches(4.4),
                Inches(0.55),
                Inches(0.45),
                "→",
                20,
                True,
                RED,
                PP_ALIGN.CENTER,
            )

    kpis = [("3,36 млн", "working"), ("18,7k", "взносы Y1"), ("781 млн ₽", "AUM Y3")]
    kw = Inches(3.95)
    for i, (v, l) in enumerate(kpis):
        x = Inches(0.4) + i * (kw + gap)
        add_rect(slide, x, Inches(5.3), kw, Inches(1.0), BG, LINE)
        add_textbox(slide, x, Inches(5.4), kw, Inches(0.4), v, 18, True, INK, PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.85), kw, Inches(0.3), l, 11, False, MUTED, PP_ALIGN.CENTER)

    add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_market(prs)
    slide_competitors(prs)
    slide_personalization(prs)
    slide_finance(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK → {OUT}")


if __name__ == "__main__":
    main()
